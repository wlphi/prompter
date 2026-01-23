"""
Voice-Activated Teleprompter Server

FastAPI server that:
1. Receives audio from browser microphone via WebSocket
2. Performs offline speech recognition using Vosk
3. Sends recognized text back to the browser
4. Serves the teleprompter frontend
"""

import asyncio
import json
import os
import secrets
import string
from contextlib import asynccontextmanager
from datetime import datetime
from pathlib import Path

from fastapi import FastAPI, WebSocket, WebSocketDisconnect, HTTPException, UploadFile, File
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pydantic import BaseModel
from vosk import Model, KaldiRecognizer, SetLogLevel
from pptx import Presentation
from dotenv import load_dotenv

from logging_config import setup_logging, get_logger

# Load environment variables from .env file
load_dotenv()

# Setup logging
LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO")
try:
    setup_logging(LOG_LEVEL)
    logger = get_logger("main")
    logger.info(f"Logging initialized at {LOG_LEVEL} level")
except ValueError as e:
    # Fall back to INFO if invalid log level
    setup_logging("INFO")
    logger = get_logger("main")
    logger.warning(f"Invalid LOG_LEVEL '{LOG_LEVEL}', using INFO: {e}")

# Reduce Vosk logging verbosity
SetLogLevel(-1)
logger.debug("Vosk logging disabled")

# Configuration
MODELS_DIR = Path(__file__).parent / "models"
SCRIPTS_DIR = Path(__file__).parent / "scripts"
FRONTEND_DIR = Path(__file__).parent.parent / "frontend"
SAMPLE_RATE = 16000
CODE_LENGTH = 5
EXPIRE_DAYS = 30

logger.info(f"Configuration: models={MODELS_DIR}, scripts={SCRIPTS_DIR}, frontend={FRONTEND_DIR}")

# Ensure scripts directory exists
SCRIPTS_DIR.mkdir(exist_ok=True)
logger.debug(f"Scripts directory ready: {SCRIPTS_DIR}")

# Cache loaded models
loaded_models: dict[str, Model] = {}


def get_model(model_path: str) -> Model:
    """Get or load a Vosk model."""
    if model_path not in loaded_models:
        logger.info(f"Loading Vosk model: {model_path}")
        try:
            loaded_models[model_path] = Model(model_path)
            logger.info(f"Model loaded successfully: {model_path}")
        except Exception as e:
            logger.error(f"Failed to load model {model_path}: {e}")
            raise
    else:
        logger.debug(f"Using cached model: {model_path}")
    return loaded_models[model_path]


def get_available_models() -> list[dict]:
    """List available Vosk models in the models directory."""
    logger.debug(f"Scanning for models in: {MODELS_DIR}")
    models = []
    if MODELS_DIR.exists():
        for model_dir in MODELS_DIR.iterdir():
            if model_dir.is_dir() and (model_dir / "am" / "final.mdl").exists():
                name = model_dir.name
                parts = name.replace("vosk-model-", "").replace("small-", "").replace("big-", "").split("-")
                if len(parts) >= 2:
                    lang_code = f"{parts[0]}-{parts[1]}" if len(parts[1]) == 2 else parts[0]
                else:
                    lang_code = parts[0] if parts else "unknown"
                models.append({
                    "path": str(model_dir),
                    "name": name,
                    "language": lang_code,
                })
                logger.debug(f"Found model: {name} ({lang_code})")
    else:
        logger.warning(f"Models directory does not exist: {MODELS_DIR}")

    logger.info(f"Found {len(models)} available models")
    return models


# Script storage
class ScriptData(BaseModel):
    script: str
    language: str = ""
    fontSize: int = 48
    scrollMargin: int = 30


def generate_code() -> str:
    """Generate a random 5-character code."""
    chars = string.ascii_lowercase + string.digits
    while True:
        code = ''.join(secrets.choice(chars) for _ in range(CODE_LENGTH))
        if not (SCRIPTS_DIR / f"{code}.json").exists():
            return code


def save_script(code: str, data: ScriptData) -> None:
    """Save script data to file."""
    now = datetime.utcnow().isoformat() + "Z"
    script_file = SCRIPTS_DIR / f"{code}.json"

    # Preserve created date if updating
    created = now
    if script_file.exists():
        existing = json.loads(script_file.read_text())
        created = existing.get("created", now)

    script_file.write_text(json.dumps({
        "script": data.script,
        "language": data.language,
        "fontSize": data.fontSize,
        "scrollMargin": data.scrollMargin,
        "created": created,
        "accessed": now
    }, indent=2))


def load_script(code: str) -> dict | None:
    """Load script data from file."""
    script_file = SCRIPTS_DIR / f"{code}.json"
    if not script_file.exists():
        return None

    data = json.loads(script_file.read_text())

    # Update access time
    data["accessed"] = datetime.utcnow().isoformat() + "Z"
    script_file.write_text(json.dumps(data, indent=2))

    return data


@asynccontextmanager
async def lifespan(app: FastAPI):
    """Application lifespan."""
    # Startup
    logger.info("=== Teleprompter Server Starting ===")
    logger.info(f"Frontend directory: {FRONTEND_DIR} (exists: {FRONTEND_DIR.exists()})")
    logger.info(f"Scripts directory: {SCRIPTS_DIR} (exists: {SCRIPTS_DIR.exists()})")
    logger.info(f"Models directory: {MODELS_DIR} (exists: {MODELS_DIR.exists()})")

    # Log available models
    models = get_available_models()
    if models:
        logger.info(f"Available models: {', '.join(m['language'] for m in models)}")
    else:
        logger.warning("No models available - Vosk recognition will not work")

    logger.info("=== Server Ready ===")

    yield

    # Shutdown
    logger.info("=== Server Shutting Down ===")
    if loaded_models:
        logger.info(f"Unloading {len(loaded_models)} cached models")
    logger.info("=== Server Stopped ===")


app = FastAPI(
    title="Voice-Activated Teleprompter",
    description="Offline speech recognition teleprompter",
    lifespan=lifespan
)


@app.get("/api/models")
async def api_get_models():
    """List available speech recognition models."""
    logger.debug("API: Listing available models")
    models = get_available_models()
    return {"models": models}


@app.post("/api/scripts")
async def api_save_script(data: ScriptData):
    """Save a script and return the access code."""
    logger.info(f"API: Saving script (length: {len(data.script)} chars)")
    code = generate_code()
    save_script(code, data)
    logger.info(f"API: Script saved with code: {code}")
    return {"code": code}


@app.get("/api/scripts/{code}")
async def api_load_script(code: str):
    """Load a script by its code."""
    logger.info(f"API: Loading script with code: {code}")

    # Sanitize code
    code = code.lower().strip()[:CODE_LENGTH]
    if not code.isalnum():
        logger.warning(f"API: Invalid code format: {code}")
        raise HTTPException(status_code=400, detail="Invalid code")

    data = load_script(code)
    if data is None:
        logger.warning(f"API: Script not found: {code}")
        raise HTTPException(status_code=404, detail="Script not found")

    logger.info(f"API: Script loaded successfully: {code}")
    return data


@app.post("/api/import/pptx")
async def api_import_pptx(file: UploadFile = File(...)):
    """Extract speaker notes from a PowerPoint file."""
    logger.info(f"API: Importing PPTX: {file.filename}")

    if not file.filename.lower().endswith('.pptx'):
        logger.warning(f"API: Invalid file type: {file.filename}")
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    try:
        contents = await file.read()
        file_size = len(contents)
        logger.debug(f"API: PPTX file size: {file_size} bytes")

        from io import BytesIO
        prs = Presentation(BytesIO(contents))

        parts = []
        for i, slide in enumerate(prs.slides):
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(notes_text)

            # Add slide transition cue (except after last slide)
            if i < len(prs.slides) - 1:
                parts.append("[next slide]")

        script = "\n\n".join(parts)
        logger.info(f"API: Extracted {len(parts)} note sections from {len(prs.slides)} slides ({len(script)} chars)")
        return {"script": script, "slideCount": len(prs.slides)}

    except Exception as e:
        logger.error(f"API: Failed to parse PPTX {file.filename}: {e}", exc_info=True)
        raise HTTPException(status_code=400, detail=f"Failed to parse PPTX: {str(e)}")


@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    """
    WebSocket endpoint for audio streaming and recognition.

    Protocol:
    - Client sends: {"type": "start", "model": "/path/to/model"}
    - Client sends: binary audio data (16-bit PCM, 16kHz, mono)
    - Client sends: {"type": "stop"}
    - Server sends: {"type": "partial", "text": "...", "words": [...]}
    - Server sends: {"type": "final", "text": "...", "words": [...]}
    """
    # Generate unique client ID for tracking
    client_id = id(websocket)
    client_host = websocket.client.host if websocket.client else "unknown"
    logger.info(f"[WS:{client_id}] Client connecting from {client_host}")

    await websocket.accept()
    logger.debug(f"[WS:{client_id}] Connection accepted")

    recognizer: KaldiRecognizer | None = None
    audio_chunks_received = 0

    try:
        while True:
            message = await websocket.receive()

            if "text" in message:
                # JSON control message
                data = json.loads(message["text"])
                msg_type = data.get("type")
                logger.debug(f"[WS:{client_id}] Received control message: {msg_type}")

                if msg_type == "start":
                    model_path = data.get("model")
                    logger.info(f"[WS:{client_id}] Starting recognition with model: {model_path}")

                    if not model_path or not Path(model_path).exists():
                        logger.warning(f"[WS:{client_id}] Invalid model path: {model_path}")
                        await websocket.send_text(json.dumps({
                            "type": "error",
                            "message": "Invalid model path"
                        }))
                        continue

                    try:
                        model = get_model(model_path)
                        recognizer = KaldiRecognizer(model, SAMPLE_RATE)
                        recognizer.SetWords(True)
                        audio_chunks_received = 0
                        logger.info(f"[WS:{client_id}] Recognizer initialized successfully")
                        await websocket.send_text(json.dumps({
                            "type": "ready"
                        }))
                    except Exception as e:
                        logger.error(f"[WS:{client_id}] Failed to initialize recognizer: {e}", exc_info=True)
                        await websocket.send_text(json.dumps({
                            "type": "error",
                            "message": str(e)
                        }))

                elif msg_type == "stop":
                    logger.info(f"[WS:{client_id}] Stop signal received (processed {audio_chunks_received} audio chunks)")
                    if recognizer:
                        # Get final result
                        result = json.loads(recognizer.FinalResult())
                        text = result.get("text", "").strip()
                        if text:
                            logger.debug(f"[WS:{client_id}] Final result on stop: '{text}'")
                            await websocket.send_text(json.dumps({
                                "type": "final",
                                "text": text,
                                "words": text.lower().split()
                            }))
                    recognizer = None
                    await websocket.send_text(json.dumps({
                        "type": "stopped"
                    }))

                elif msg_type == "ping":
                    logger.debug(f"[WS:{client_id}] Ping received, sending pong")
                    await websocket.send_text(json.dumps({"type": "pong"}))

            elif "bytes" in message:
                # Binary audio data
                if recognizer is None:
                    logger.debug(f"[WS:{client_id}] Received audio but recognizer not initialized, ignoring")
                    continue

                audio_data = message["bytes"]
                audio_chunks_received += 1
                logger.debug(f"[WS:{client_id}] Processing audio chunk {audio_chunks_received} ({len(audio_data)} bytes)")

                if recognizer.AcceptWaveform(audio_data):
                    result = json.loads(recognizer.Result())
                    text = result.get("text", "").strip()
                    if text:
                        logger.info(f"[WS:{client_id}] Final result: '{text}'")
                        await websocket.send_text(json.dumps({
                            "type": "final",
                            "text": text,
                            "words": text.lower().split()
                        }))
                else:
                    partial = json.loads(recognizer.PartialResult())
                    text = partial.get("partial", "").strip()
                    if text:
                        logger.debug(f"[WS:{client_id}] Partial result: '{text}'")
                        await websocket.send_text(json.dumps({
                            "type": "partial",
                            "text": text,
                            "words": text.lower().split()
                        }))

    except WebSocketDisconnect:
        logger.info(f"[WS:{client_id}] Client disconnected normally (processed {audio_chunks_received} audio chunks)")
    except Exception as e:
        logger.error(f"[WS:{client_id}] WebSocket error: {e}", exc_info=True)
        try:
            await websocket.send_text(json.dumps({
                "type": "error",
                "message": str(e)
            }))
        except:
            logger.debug(f"[WS:{client_id}] Could not send error to disconnected client")


@app.get("/")
async def serve_index():
    """Serve the main page."""
    index_path = FRONTEND_DIR / "index.html"
    if index_path.exists():
        return FileResponse(index_path)
    return {"error": "Frontend not found"}


# Mount static files (CSS, JS)
if FRONTEND_DIR.exists():
    app.mount("/static", StaticFiles(directory=str(FRONTEND_DIR)), name="static")


if __name__ == "__main__":
    import uvicorn

    # Get configuration from environment
    host = os.getenv("HOST", "0.0.0.0")
    port = int(os.getenv("PORT", "8000"))

    logger.info(f"Starting uvicorn server on {host}:{port}")

    uvicorn.run(
        app,
        host=host,
        port=port,
        log_level=LOG_LEVEL.lower()
    )
